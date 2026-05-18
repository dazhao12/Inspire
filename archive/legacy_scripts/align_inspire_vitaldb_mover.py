#!/usr/bin/env python3
"""Create variable-alignment audit (INSPIRE vs VitalDB vs MOVER) from DOCX standard.

Inputs:
- DOCX variable definition file (INSPIRE standard)
- Data root containing INSPIRE_1.3 / VitalDB_1.0.0 / MOVER

Outputs:
- CSV detail table with per-variable availability and evidence
- CSV summary by category and database
- PPTX with check/cross tables by category
"""

import argparse
import csv
import re
from collections import defaultdict
from datetime import datetime
from pathlib import Path
from typing import DefaultDict, Dict, Iterable, List, Sequence, Set, Tuple

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Inches, Pt


REQUESTED_CATEGORIES: List[str] = [
    "基本信息",
    "术前化验",
    "合并症",
    "手术信息",
    "术中生命体征",
    "术中用药",
    "住院用药",
    "术中化验",
    "术后化验",
    "并发症",
]

DB_KEYS = ["INSPIRE", "VitalDB", "MOVER", "TJCH"]
DB_LABELS = {
    "INSPIRE": "INSPIRE",
    "VitalDB": "VitalDB",
    "MOVER": "MOVER",
    "TJCH": "TJCH",
}

CATEGORY_SLUG = {
    "基本信息": "basic_info",
    "术前化验": "preop_labs",
    "合并症": "comorbidities",
    "手术信息": "surgery_info",
    "术中生命体征": "intraop_vitals",
    "术中用药": "intraop_meds",
    "住院用药": "inhospital_meds",
    "术中化验": "intraop_labs",
    "术后化验": "postop_labs",
    "并发症": "complications",
}

MODULE_ALIASES = {
    "术前用药": "住院用药",
}

CATEGORY_DOMAIN = {
    "基本信息": "headers",
    "术前化验": "labs",
    "合并症": "dx",
    "手术信息": "headers",
    "术中生命体征": "headers",
    "术中用药": "meds",
    "住院用药": "meds",
    "术中化验": "labs",
    "术后化验": "labs",
    "并发症": "complications",
}

COMPOSITE_COMPONENTS = [
    ("死亡", ["In-hospital mortality", "30-day mortality"]),
    ("重大心血管事件", ["Acute myocardial infarction", "Cardiac arrest", "Pulmonary embolism"]),
    ("重大神经事件", ["Cerebral infarction", "Intracerebral hemorrhage", "Subarachnoid hemorrhage"]),
    ("呼吸衰竭相关", ["ARDS", "Postoperative mechanical ventilation"]),
    ("肾脏不良事件", ["Acute kidney failure (ICD-10)", "AKI by creatinine (postoperative 7-day, non-imputed)", "Postoperative CRRT use"]),
    ("感染/休克相关", ["Sepsis", "Shock", "DIC"]),
    ("治疗升级", ["Unexpected ICU admission from general ward", "Reoperation", "Postoperative ECMO use", "Postoperative IABP use"]),
]

REPORT_COLORS = {
    "navy": RGBColor(27, 46, 84),
    "blue": RGBColor(53, 94, 174),
    "light_blue": RGBColor(228, 236, 247),
    "header_gray": RGBColor(233, 236, 241),
    "text": RGBColor(35, 42, 52),
    "ok": RGBColor(34, 139, 34),
    "no": RGBColor(196, 61, 61),
}

INTRAOP_MED_KEYWORDS = {
    "dose",
    "infusion",
    "bolus",
    "concentration",
    "epinephrine",
    "norepinephrine",
    "phenylephrine",
    "dopamine",
    "dobutamine",
    "vasopressin",
    "propofol",
    "remifentanil",
    "midazolam",
    "fentanyl",
    "sufentanil",
}

STOP_TOKENS = {
    "intraoperative",
    "preoperative",
    "postoperative",
    "value",
    "values",
    "within",
    "before",
    "after",
    "total",
    "time",
    "interval",
    "duration",
    "entry",
    "exit",
    "surgery",
    "operation",
    "monitor",
    "monitoring",
    "min",
    "minute",
    "minutes",
    "day",
    "days",
    "hour",
    "hours",
    "rate",
    "dose",
    "units",
    "unit",
    "volume",
    "level",
    "panel",
    "index",
    "support",
    "exposure",
    "measured",
    "end",
    "tidal",
    "concentration",
    "target",
    "bolus",
}

WEAK_MATCH_TOKENS = {
    "other",
    "new",
    "onset",
    "history",
    "disease",
    "event",
    "events",
    "status",
    "time",
    "days",
    "day",
    "interval",
    "from",
    "to",
    "in",
    "hospital",
    "post",
}

TOKEN_ALIAS = {
    "male": {"sex", "gender"},
    "sex": {"gender", "male", "female"},
    "bmi": {"body", "mass", "index"},
    "hemoglobin": {"hgb", "hb"},
    "hb": {"hgb", "hemoglobin"},
    "hgb": {"hb", "hemoglobin"},
    "creatinine": {"cr"},
    "cr": {"creatinine"},
    "sodium": {"na"},
    "na": {"sodium"},
    "potassium": {"k"},
    "k": {"potassium"},
    "glucose": {"gluc"},
    "gluc": {"glucose"},
    "albumin": {"alb"},
    "alb": {"albumin"},
    "platelet": {"plt"},
    "plt": {"platelet"},
    "alp": {"alkaline", "phosphatase"},
    "alt": {"alanine", "aminotransferase"},
    "ast": {"aspartate", "aminotransferase"},
    "aptt": {"ptt"},
    "ptinr": {"pt", "inr"},
    "bilirubin": {"tbil"},
    "protein": {"tp"},
    "troponin": {"tni", "tnt"},
    "spo2": {"sp02", "sao2", "ssto2"},
    "sp02": {"spo2"},
    "sbp": {"nsbp", "nibp_sbp", "art_sbp", "sbp_art", "sbp_fem"},
    "dbp": {"ndbp", "nibp_dbp", "art_dbp", "dbp_art", "dbp_fem"},
    "mbp": {"map", "nmap", "nibp_mbp", "art_mbp", "mbp_merged"},
    "hr": {"hre", "hrp"},
    "rr": set(),
    "bt": {"temp"},
    "temperature": {"bt", "temp"},
    "arb": {"arbs"},
    "arbs": {"arb"},
    "antihypertensive": {"antihypertensives"},
    "antihypertensives": {"antihypertensive"},
    "immunoglobulin": {"ivig"},
    "ivig": {"immunoglobulin"},
    "mortality": {"death", "expired"},
    "death": {"mortality"},
    "reoperation": {"reop"},
    "ventilation": {"vent", "ventilator"},
    "icu": {"intensive", "care"},
    "sepsis": {"septic"},
    "infarction": {"infarct", "mi"},
    "hemorrhage": {"bleed", "bleeding"},
    # Intraoperative medication/anesthetic abbreviations in INSPIRE/VitalDB time-series.
    "epinephrine": {"epi", "epii"},
    "norepinephrine": {"nepi", "noradrenaline", "ne"},
    "phenylephrine": {"pepi", "neo", "phen"},
    "dopamine": {"dopai", "dopa"},
    "dobutamine": {"dobui", "dobu"},
    "vasopressin": {"vaso"},
    "propofol": {"ppf", "ppfi", "ppf20"},
    "remifentanil": {"rfti", "rftn20"},
    "fentanyl": {"ftn"},
    "midazolam": {"mdz"},
    "sufentanil": {"sft"},
    "sevoflurane": {"etsevo", "exp_sevo", "insp_sevo", "sevo"},
    "desflurane": {"etdes", "exp_des", "insp_des", "des"},
    "isoflurane": {"etiso", "iso"},
    "anesthetic": {"etgas", "mac"},
    "gas": {"etgas", "mac"},
}

TOKEN_ALIAS_PHRASE = {
    "sbp": ["systolic blood pressure"],
    "dbp": ["diastolic blood pressure"],
    "mbp": ["mean blood pressure"],
    "hr": ["heart rate"],
    "rr": ["respiratory rate"],
    "bt": ["body temperature"],
    "spo2": ["oxygen saturation"],
}

# Explicitly block known wrong mappings found during manual review.
# key: (normalized variable name, db_key), value: blocked normalized patterns.
FORBIDDEN_MATCH_PATTERNS = {
    ("preoperative spo2", "TJCH"): ["ssto2"],
    ("preoperative chloride", "TJCH"): ["mivacurium chloride"],
    ("preoperative troponin i", "TJCH"): ["aki stages i"],
    ("beta blockers", "MOVER"): ["beta carotene"],
    ("ace inhibitors", "MOVER"): ["medroxyprogest ace", "ace pramoxine"],
}

PREOP_POSITIVE_TOKENS = [
    "preop",
    "pre op",
    "pre operative",
    "preoperative",
    "baseline",
    "admission",
    "before surgery",
    "before operation",
]

PREOP_NEGATIVE_TOKENS = [
    "postop",
    "post op",
    "post operative",
    "postoperative",
    "postdischarge",
    "intraop",
    "intra op",
    "intra operative",
    "intraoperative",
]

PREOP_TIME_CHECK_CATEGORIES = {
    "术前化验",
    "住院用药",  # "术前用药" is mapped to this module via MODULE_ALIASES.
}

COMORBIDITY_POSITIVE_TOKENS = [
    "comorbidity",
    "comorbid",
    "history",
    "patient history",
    "preop",
    "pre op",
    "preoperative",
    "baseline",
    "admission",
    "chronic",
]

# User-confirmed TJCH preoperative medication fields.
# key: standard variable name (normalized), value: TJCH header names in cohort xlsx.
TJCH_PREOP_MED_OVERRIDES: Dict[str, List[str]] = {
    "beta blockers": ["Beta_blockers"],
    "ace inhibitors": ["ACEI"],
    "calcium channel blockers": ["CCB"],
    "nitrates": ["Nitroglycerin"],
    "diuretics": ["Diuretics"],
    "antidiabetic drugs": ["Antidiabetic_agents", "Oral_antidiabetic"],
    "insulins": ["Insulin_injection", "Insulin"],
    "arbs": ["ARB"],
    "statins": ["Statins"],
    "anticoagulants": ["Anticoagulation"],
}

# Variable-level medication synonyms (applied only for 住院用药 matching).
MED_VAR_SYNONYM_PHRASES: Dict[str, List[str]] = {
    "beta blockers": ["metoprolol", "atenolol", "carvedilol", "bisoprolol", "propranolol", "labetalol", "esmolol"],
    "ace inhibitors": ["lisinopril", "enalapril", "ramipril", "benazepril", "captopril", "fosinopril", "quinapril", "perindopril"],
    "calcium channel blockers": ["amlodipine", "diltiazem", "verapamil", "nifedipine", "nicardipine", "clevidipine", "felodipine"],
    "arbs": ["losartan", "valsartan", "irbesartan", "candesartan", "telmisartan", "olmesartan", "azilsartan"],
    "nitrates": ["nitroglycerin", "isosorbide", "nitroprusside"],
    "diuretics": ["furosemide", "bumetanide", "torsemide", "spironolactone", "hydrochlorothiazide", "chlorothiazide", "metolazone"],
    "anticoagulants": ["heparin", "enoxaparin", "warfarin", "apixaban", "rivaroxaban", "dabigatran", "fondaparinux", "argatroban"],
    "statins": ["atorvastatin", "rosuvastatin", "simvastatin", "pravastatin", "lovastatin", "fluvastatin", "pitavastatin"],
    "insulins": ["insulin lispro", "insulin aspart", "insulin glargine", "insulin detemir", "insulin regular", "insulin nph", "insulin"],
    "antidiabetic drugs": [
        "metformin",
        "glipizide",
        "glyburide",
        "pioglitazone",
        "sitagliptin",
        "linagliptin",
        "empagliflozin",
        "dapagliflozin",
        "canagliflozin",
        "liraglutide",
        "semaglutide",
    ],
    # INSPIRE preop_meds_defined naming harmonization
    "other antihypertensives": ["other antihypertensive"],
    "intravenous immunoglobulin": ["ivig"],
    "anti tuberculosis drugs": ["anti tuberculosis", "anti_tuberculosis"],
    "5 ht3 receptor antagonists": ["5ht3 antagonists", "serotonin 5ht3 antagonists", "serotonin_5ht3_antagonists"],
}


def variable_phase(var_name: str) -> str:
    nv = normalize_text(var_name)
    if nv.startswith("preoperative "):
        return "preop"
    if nv.startswith("postoperative "):
        return "postop"
    if nv.startswith("intraoperative "):
        return "intraop"
    return ""


def source_phase(src_norm: str) -> str:
    s = src_norm
    if any(k in s for k in ["preop", "pre op", "preoperative", "baseline"]):
        return "preop"
    if any(k in s for k in ["postop", "post op", "postoperative", "postdischarge"]):
        return "postop"
    if any(
        k in s
        for k in [
            "intraop",
            "intra op",
            "intraoperative",
            "solar8000",
            "primus",
            "orchestra",
            "vigilance",
            "cardioq",
            "ev1000",
            "fms",
            "bis",
            "invos",
            "art sbp",
            "nibp sbp",
            "art dbp",
            "nibp dbp",
            "pap",
            "cvp",
            "etco2",
            "fio2",
            "peep",
            "pip",
            "ppf",
            "rftn",
            "intra1",
            "intra2",
            " intra ",
        ]
    ):
        return "intraop"
    return ""


def override_match(var_name: str, db_key: str, catalog) -> Tuple[bool, List[str]]:
    nv = normalize_text(var_name)
    if db_key == "TJCH":
        manual = TJCH_PREOP_MED_OVERRIDES.get(nv, [])
        if manual:
            hits: List[str] = []
            for m in manual:
                nm = normalize_text(m)
                if nm in catalog.norm_to_raw:
                    hits.extend(sorted(catalog.norm_to_raw[nm]))
            if hits:
                # Keep deterministic order and avoid duplicates.
                uniq: List[str] = []
                seen: Set[str] = set()
                for h in hits:
                    if h in seen:
                        continue
                    seen.add(h)
                    uniq.append(h)
                return True, uniq[:3]

    patterns: List[str] = []

    if "aki by creatinine" in nv:
        patterns = [" aki ", "aki ", " aki", "cr highest 7d", "cr 0 7d high", "peak creatinine"]
    elif "acute kidney failure icd 10" in nv:
        patterns = ["acute kidney failure", "aki", "aki cat"]
    elif "postoperative peak creatinine" in nv:
        patterns = ["peak creatinine", "cr highest 7d", "cr 0 7d high"]

    if not patterns:
        return False, []

    hits = []
    for src in catalog.norm_sources:
        s = " " + src + " "
        if any(p in s for p in patterns):
            hits.extend(sorted(catalog.norm_to_raw[src]))
            if len(hits) >= 3:
                break
    if hits:
        return True, hits[:3]
    return False, []


def normalize_text(text: str) -> str:
    if text is None:
        return ""
    s = str(text).strip().lower()
    s = s.replace("₂", "2")
    s = s.replace("₃", "3")
    s = s.replace("₅", "5")
    s = s.replace("₇", "7")
    s = s.replace("₉", "9")
    s = s.replace("₀", "0")
    s = s.replace("%", " percent ")
    s = s.replace("°", " degree ")
    s = re.sub(r"[^a-z0-9]+", " ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s


def tokenize(text: str) -> Set[str]:
    toks = set(normalize_text(text).split())
    return {t for t in toks if t and t not in STOP_TOKENS}


def remove_context_words(norm_var: str) -> str:
    context_words = {
        "preoperative",
        "postoperative",
        "intraoperative",
        "admission",
        "postoperative",
        "or",
        "entry",
        "exit",
    }
    toks = [t for t in norm_var.split() if t not in context_words]
    return " ".join(toks).strip()


def iter_table_variables(doc: Document, table_idx: int) -> List[str]:
    t = doc.tables[table_idx]
    out: List[str] = []
    for row in t.rows[1:]:
        vals = [c.text.strip() for c in row.cells]
        if len(set(vals)) == 1:
            continue
        var = vals[0].strip()
        if var:
            out.append(var)
    return out


def split_intraop_table(vars_intraop: Sequence[str]) -> Tuple[List[str], List[str]]:
    vitals: List[str] = []
    meds: List[str] = []
    for v in vars_intraop:
        nv = normalize_text(v)
        if any(k in nv for k in INTRAOP_MED_KEYWORDS):
            meds.append(v)
        else:
            vitals.append(v)
    return vitals, meds


def extract_standard_variables(docx_path: Path) -> Dict[str, List[str]]:
    doc = Document(str(docx_path))

    t0 = iter_table_variables(doc, 0)  # demographics/baseline
    t1 = iter_table_variables(doc, 1)  # comorbidities
    t2 = iter_table_variables(doc, 2)  # baseline vitals
    t3 = iter_table_variables(doc, 3)  # preop labs
    t4 = iter_table_variables(doc, 4)  # medications
    t5 = iter_table_variables(doc, 5)  # timeline
    t6 = iter_table_variables(doc, 6)  # intraop management
    t7 = iter_table_variables(doc, 7)  # postop labs
    t8 = iter_table_variables(doc, 8)  # outcomes/complications

    intraop_vitals, intraop_meds = split_intraop_table(t6)

    # DOCX has no standalone intraop-lab table; derive analyte names from preop labs.
    intraop_labs = []
    for v in t3:
        nv = normalize_text(v)
        if nv.startswith("preoperative "):
            intraop_labs.append(v.replace("Preoperative", "Intraoperative", 1))
        else:
            intraop_labs.append(f"Intraoperative {v}")

    categories: Dict[str, List[str]] = {
        "基本信息": t0 + t2,
        "术前化验": t3,
        "合并症": t1,
        "手术信息": t5,
        "术中生命体征": intraop_vitals,
        "术中用药": intraop_meds,
        "住院用药": t4,
        "术中化验": intraop_labs,
        "术后化验": t7,
        "并发症": t8,
    }

    for k in REQUESTED_CATEGORIES:
        categories.setdefault(k, [])

    return categories


def parse_modules_arg(modules_arg: str) -> List[str]:
    if not modules_arg:
        return list(REQUESTED_CATEGORIES)
    raw = [x.strip() for x in modules_arg.split(",") if x.strip()]
    if not raw:
        return list(REQUESTED_CATEGORIES)
    selected: List[str] = []
    unknown: List[str] = []
    for m in raw:
        m = MODULE_ALIASES.get(m, m)
        if m in REQUESTED_CATEGORIES:
            selected.append(m)
        else:
            unknown.append(m)
    if unknown:
        raise ValueError("Unknown modules: " + ", ".join(unknown))
    return selected


class DBCatalog:
    def __init__(self, name, headers, values, header_sources=None, value_sources=None):
        self.name = name
        self.headers = set(headers)
        self.values = set(values)
        self.header_sources = header_sources or {}
        self.value_sources = value_sources or {}

        raw = set(self.headers) | set(self.values)
        self.raw_to_sources: DefaultDict[str, Set[str]] = defaultdict(set)
        for h in self.headers:
            self.raw_to_sources[h].update(self.header_sources.get(h, set()))
        for v in self.values:
            self.raw_to_sources[v].update(self.value_sources.get(v, set()))
        for x in raw:
            if not self.raw_to_sources[x]:
                self.raw_to_sources[x].add("unknown_source")

        self.norm_to_raw: Dict[str, Set[str]] = defaultdict(set)
        for x in raw:
            nx = normalize_text(x)
            if nx:
                self.norm_to_raw[nx].add(x)
        self.norm_sources: List[str] = sorted(self.norm_to_raw.keys())
        self.source_tokens: List[Set[str]] = [tokenize(x) for x in self.norm_sources]
        self.token_index: Dict[str, Set[int]] = defaultdict(set)
        for i, toks in enumerate(self.source_tokens):
            for t in toks:
                self.token_index[t].add(i)

    def sources_for_raw(self, raw_values: Sequence[str], limit: int = 6) -> List[str]:
        srcs: Set[str] = set()
        for raw in raw_values:
            srcs.update(self.raw_to_sources.get(raw, set()))
        return sorted(srcs)[:limit]


def get_header(csv_path: Path) -> List[str]:
    try:
        with csv_path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
            row = next(csv.reader(f), [])
            return [x.strip() for x in row if x and x.strip()]
    except Exception:
        return []


def merge_source_maps(dst: DefaultDict[str, Set[str]], src: Dict[str, Set[str]]) -> None:
    for key, vals in src.items():
        dst[key].update(vals)


def scan_csv_headers_with_sources(csv_files: Iterable[Path], root: Path) -> Tuple[Set[str], Dict[str, Set[str]]]:
    headers: Set[str] = set()
    source_map: DefaultDict[str, Set[str]] = defaultdict(set)
    for p in csv_files:
        rel = str(p.relative_to(root)) if root in p.parents else str(p)
        for h in get_header(p):
            headers.add(h)
            source_map[h].add(rel)
    return headers, dict(source_map)


def get_excel_headers(xlsx_path: Path) -> Set[str]:
    headers: Set[str] = set()
    try:
        wb = load_workbook(filename=str(xlsx_path), read_only=True, data_only=True)
    except Exception:
        return headers

    for ws in wb.worksheets:
        try:
            first = next(ws.iter_rows(min_row=1, max_row=1, values_only=True))
        except Exception:
            continue
        for val in first:
            if val is None:
                continue
            name = str(val).strip()
            if name:
                headers.add(name)
    wb.close()
    return headers


def get_excel_headers_with_sources(xlsx_path: Path, root: Path) -> Tuple[Set[str], Dict[str, Set[str]]]:
    headers: Set[str] = set()
    source_map: DefaultDict[str, Set[str]] = defaultdict(set)
    rel = str(xlsx_path.relative_to(root)) if root in xlsx_path.parents else str(xlsx_path)
    try:
        wb = load_workbook(filename=str(xlsx_path), read_only=True, data_only=True)
    except Exception:
        return headers, {}

    for ws in wb.worksheets:
        try:
            first = next(ws.iter_rows(min_row=1, max_row=1, values_only=True))
        except Exception:
            continue
        for val in first:
            if val is None:
                continue
            name = str(val).strip()
            if name:
                headers.add(name)
                source_map[name].add(f"{rel}#{ws.title}")
    wb.close()
    return headers, dict(source_map)


def get_headers_from_csv_dir(csv_dir: Path, pattern: str = "*.csv", max_files: int = 0) -> Set[str]:
    headers: Set[str] = set()
    if not csv_dir.exists():
        return headers
    files = sorted(csv_dir.glob(pattern))
    if max_files and max_files > 0:
        files = files[:max_files]
    for p in files:
        headers.update(get_header(p))
    return headers


def collect_unique_values_with_sources(
    csv_path: Path,
    cols: Sequence[str],
    root: Path,
    limit: int = 20000,
    max_rows: int = 300000,
) -> Tuple[Set[str], Dict[str, Set[str]]]:
    out: Set[str] = set()
    source_map: DefaultDict[str, Set[str]] = defaultdict(set)
    if not cols:
        return out, {}

    normalized_cols = {c.lower(): c for c in cols}
    rel = str(csv_path.relative_to(root)) if root in csv_path.parents else str(csv_path)
    try:
        with csv_path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.DictReader(f)
            if not reader.fieldnames:
                return out, {}
            keymap = {}
            for fn in reader.fieldnames:
                if fn is None:
                    continue
                low = fn.lower().strip()
                if low in normalized_cols:
                    keymap[fn] = normalized_cols[low]
            if not keymap:
                return out, {}
            row_count = 0
            for row in reader:
                row_count += 1
                for src_key in keymap:
                    val = row.get(src_key)
                    if val is None:
                        continue
                    sval = str(val).strip()
                    if not sval:
                        continue
                    out.add(sval)
                    source_map[sval].add(f"{rel}:{src_key}")
                    if len(out) >= limit:
                        return out, dict(source_map)
                if row_count >= max_rows:
                    return out, dict(source_map)
    except Exception:
        return out, dict(source_map)

    return out, dict(source_map)


def build_catalogs(data_root: Path, tjth_xlsx: Path) -> Dict[str, Dict[str, DBCatalog]]:
    catalogs: Dict[str, Dict[str, DBCatalog]] = {}

    # INSPIRE: scan full tree headers (manageable size).
    inspire_root = data_root / "INSPIRE_1.3"
    inspire_headers, inspire_header_sources = scan_csv_headers_with_sources(inspire_root.rglob("*.csv"), data_root)
    inspire_med_values: Set[str] = set()
    inspire_med_sources: DefaultDict[str, Set[str]] = defaultdict(set)
    inspire_param = inspire_root / "01_source_raw" / "parameters.csv"
    inspire_vitals = inspire_root / "01_source_raw" / "vitals.csv"
    if inspire_param.exists():
        vals, srcs = collect_unique_values_with_sources(
            inspire_param,
            ["Label"],
            root=data_root,
            limit=30000,
            max_rows=500000,
        )
        inspire_med_values.update(vals)
        merge_source_maps(inspire_med_sources, srcs)
    if inspire_vitals.exists():
        vals, srcs = collect_unique_values_with_sources(
            inspire_vitals,
            ["item_name"],
            root=data_root,
            limit=20000,
            max_rows=3500000,
        )
        inspire_med_values.update(vals)
        merge_source_maps(inspire_med_sources, srcs)
    catalogs["INSPIRE"] = {
        "headers": DBCatalog("INSPIRE_headers", headers=inspire_headers, values=set(), header_sources=inspire_header_sources),
        "labs": DBCatalog("INSPIRE_labs", headers=inspire_headers, values=set(), header_sources=inspire_header_sources),
        "meds": DBCatalog(
            "INSPIRE_meds",
            headers=inspire_headers,
            values=inspire_med_values,
            header_sources=inspire_header_sources,
            value_sources=dict(inspire_med_sources),
        ),
        "dx": DBCatalog("INSPIRE_dx", headers=inspire_headers, values=set(), header_sources=inspire_header_sources),
        "complications": DBCatalog("INSPIRE_comp", headers=inspire_headers, values=set(), header_sources=inspire_header_sources),
    }

    # VitalDB: focus on processed tables to avoid per-case raw expansion.
    vital_root = data_root / "VitalDB_1.0.0"
    vital_headers, vital_header_sources = scan_csv_headers_with_sources((vital_root / "processed").rglob("*.csv"), data_root)
    vital_comp_headers, vital_comp_header_sources = scan_csv_headers_with_sources(
        (vital_root / "processed" / "emr_data").glob("*.csv"),
        data_root,
    )
    vital_lab_values: Set[str] = set()
    vital_lab_sources: DefaultDict[str, Set[str]] = defaultdict(set)
    vital_med_values: Set[str] = set()
    vital_med_sources: DefaultDict[str, Set[str]] = defaultdict(set)

    vital_value_cols = {
        vital_root / "processed" / "emr_data" / "intraop_lab.csv": (["name"], 50000, 400000),
        vital_root / "processed" / "emr_data" / "postop_lab.csv": (["name"], 50000, 400000),
        vital_root / "processed" / "emr_data" / "postdischarge_lab.csv": (["name"], 50000, 400000),
    }
    for p, spec in vital_value_cols.items():
        if p.exists():
            cols, lim, max_rows = spec
            vals, srcs = collect_unique_values_with_sources(p, cols, root=data_root, limit=lim, max_rows=max_rows)
            vital_lab_values.update(vals)
            merge_source_maps(vital_lab_sources, srcs)

    # Explicitly include per-case time-series tracks from csv_output for intraop meds/anesthetics.
    csv_output_headers, csv_output_header_sources = scan_csv_headers_with_sources(
        (vital_root / "csv_output").glob("*.csv"),
        data_root,
    )
    med_track_patterns = [
        "ORCHESTRA/",
        "EPI",
        "NEPI",
        "DOPA",
        "DOBU",
        "VASO",
        "PPF",
        "RFTN",
        "ETDES",
        "ETISO",
        "ETSEVO",
        "EXP_DES",
        "EXP_SEVO",
        "INSP_DES",
        "INSP_SEVO",
        "MAC",
        "SET_AGE",
    ]
    for col in (vital_headers | csv_output_headers):
        u = col.upper()
        if any(pat in u for pat in med_track_patterns):
            vital_med_values.add(col)
            for src in vital_header_sources.get(col, set()):
                vital_med_sources[col].add(src)
            for src in csv_output_header_sources.get(col, set()):
                vital_med_sources[col].add(src)

    catalogs["VitalDB"] = {
        "headers": DBCatalog("VitalDB_headers", headers=vital_headers, values=set(), header_sources=vital_header_sources),
        "labs": DBCatalog(
            "VitalDB_labs",
            headers=vital_headers,
            values=vital_lab_values,
            header_sources=vital_header_sources,
            value_sources=dict(vital_lab_sources),
        ),
        "meds": DBCatalog(
            "VitalDB_meds",
            headers=vital_headers,
            values=vital_med_values,
            header_sources=vital_header_sources,
            value_sources=dict(vital_med_sources),
        ),
        "dx": DBCatalog("VitalDB_dx", headers=vital_headers, values=set(), header_sources=vital_header_sources),
        "complications": DBCatalog(
            "VitalDB_comp",
            headers=vital_comp_headers,
            values=set(),
            header_sources=vital_comp_header_sources,
        ),
    }

    # MOVER: use raw SIS/EPIC + derived medication tables.
    mover_root = data_root / "MOVER"
    mover_headers: Set[str] = set()
    mover_header_sources: DefaultDict[str, Set[str]] = defaultdict(set)
    mover_lab_values: Set[str] = set()
    mover_lab_sources: DefaultDict[str, Set[str]] = defaultdict(set)
    mover_med_values: Set[str] = set()
    mover_med_sources: DefaultDict[str, Set[str]] = defaultdict(set)
    mover_dx_values: Set[str] = set()
    mover_dx_sources: DefaultDict[str, Set[str]] = defaultdict(set)
    mover_comp_values: Set[str] = set()
    mover_comp_sources: DefaultDict[str, Set[str]] = defaultdict(set)

    mover_scan_dirs = [
        mover_root / "raw" / "sis_EMR",
        mover_root / "raw" / "EPIC_EMR",
        mover_root / "derived" / "patient_medications_phased",
        mover_root / "derived" / "patient_meds_summaries",
    ]
    for d in mover_scan_dirs:
        if not d.exists():
            continue
        hs, srcs = scan_csv_headers_with_sources(d.rglob("*.csv"), data_root)
        mover_headers.update(hs)
        merge_source_maps(mover_header_sources, srcs)

    mover_lab_cols = {
        mover_root / "raw" / "EPIC_EMR" / "patient_labs.csv": (["Lab Name"], 40000, 400000),
    }
    mover_med_cols = {
        mover_root / "raw" / "sis_EMR" / "patient_medication.csv": (["Drug_name"], 60000, 400000),
        mover_root / "raw" / "EPIC_EMR" / "patient_medications.csv": (["DISPLAY_NAME", "MEDICATION_NM"], 90000, 500000),
        mover_root / "derived" / "patient_medications_phased" / "patient_medications_PRE_OP.csv": (
            ["DISPLAY_NAME", "MEDICATION_NM"],
            120000,
            800000,
        ),
    }
    mover_dx_cols = {
        mover_root / "raw" / "EPIC_EMR" / "patient_history.csv": (["dx_name"], 70000, 120000),
        mover_root / "raw" / "EPIC_EMR" / "patient_visit.csv": (["dx_name"], 70000, 120000),
    }
    mover_comp_cols = {
        mover_root / "raw" / "EPIC_EMR" / "patient_post_op_complications.csv": (["Element_Name", "Element_abbr"], 50000, 500000),
        mover_root / "raw" / "EPIC_EMR" / "patient_information.csv": (["DISCH_DISP"], 1000, 400000),
    }

    for p, spec in mover_lab_cols.items():
        if p.exists():
            cols, lim, max_rows = spec
            vals, srcs = collect_unique_values_with_sources(p, cols, root=data_root, limit=lim, max_rows=max_rows)
            mover_lab_values.update(vals)
            merge_source_maps(mover_lab_sources, srcs)
    for p, spec in mover_med_cols.items():
        if p.exists():
            cols, lim, max_rows = spec
            vals, srcs = collect_unique_values_with_sources(p, cols, root=data_root, limit=lim, max_rows=max_rows)
            mover_med_values.update(vals)
            merge_source_maps(mover_med_sources, srcs)
    for p, spec in mover_dx_cols.items():
        if p.exists():
            cols, lim, max_rows = spec
            vals, srcs = collect_unique_values_with_sources(p, cols, root=data_root, limit=lim, max_rows=max_rows)
            mover_dx_values.update(vals)
            merge_source_maps(mover_dx_sources, srcs)
    for p, spec in mover_comp_cols.items():
        if p.exists():
            cols, lim, max_rows = spec
            vals, srcs = collect_unique_values_with_sources(p, cols, root=data_root, limit=lim, max_rows=max_rows)
            mover_comp_values.update(vals)
            merge_source_maps(mover_comp_sources, srcs)

    catalogs["MOVER"] = {
        "headers": DBCatalog("MOVER_headers", headers=mover_headers, values=set(), header_sources=dict(mover_header_sources)),
        "labs": DBCatalog(
            "MOVER_labs",
            headers=mover_headers,
            values=mover_lab_values,
            header_sources=dict(mover_header_sources),
            value_sources=dict(mover_lab_sources),
        ),
        "meds": DBCatalog(
            "MOVER_meds",
            headers=mover_headers,
            values=mover_med_values,
            header_sources=dict(mover_header_sources),
            value_sources=dict(mover_med_sources),
        ),
        "dx": DBCatalog(
            "MOVER_dx",
            headers=mover_headers,
            values=mover_dx_values,
            header_sources=dict(mover_header_sources),
            value_sources=dict(mover_dx_sources),
        ),
        "complications": DBCatalog(
            "MOVER_comp",
            headers=mover_headers,
            values=mover_comp_values,
            header_sources=dict(mover_header_sources),
            value_sources=dict(mover_comp_sources),
        ),
    }

    # Tianjin Chest Hospital cohort xlsx (wide table).
    tjth_headers, tjth_header_sources = get_excel_headers_with_sources(tjth_xlsx, data_root)
    catalogs["TJCH"] = {
        "headers": DBCatalog("TJCH_headers", headers=tjth_headers, values=set(), header_sources=tjth_header_sources),
        "labs": DBCatalog("TJCH_labs", headers=tjth_headers, values=set(), header_sources=tjth_header_sources),
        "meds": DBCatalog("TJCH_meds", headers=tjth_headers, values=set(), header_sources=tjth_header_sources),
        "dx": DBCatalog("TJCH_dx", headers=tjth_headers, values=set(), header_sources=tjth_header_sources),
        "complications": DBCatalog("TJCH_comp", headers=tjth_headers, values=set(), header_sources=tjth_header_sources),
    }

    return catalogs


def expanded_tokens(tokens: Set[str]) -> Set[str]:
    out = set(tokens)
    for t in list(tokens):
        out.update(TOKEN_ALIAS.get(t, set()))
    return out


def phrase_in_source(phrase: str, src: str) -> bool:
    if not phrase or not src:
        return False
    if len(phrase) < 4 and " " not in phrase:
        return False
    return bool(re.search(r"\b" + re.escape(phrase) + r"\b", src))


def derive_key_tokens(var_name: str) -> Set[str]:
    toks = tokenize(var_name)
    if "male" in toks and "sex" in toks:
        return {"sex"}
    if "bmi" in toks or {"body", "mass", "index"}.issubset(toks):
        return {"bmi"}

    filtered = {t for t in toks if t not in WEAK_MATCH_TOKENS and t != "total"}
    if len(filtered) > 1:
        filtered = {t for t in filtered if not t.isdigit()}
    if not filtered:
        filtered = set(toks)
    return filtered


def is_forbidden_candidate(var_name: str, db_key: str, src_norm: str, raw_values: Sequence[str], catalog: DBCatalog) -> bool:
    phase = variable_phase(var_name)
    texts = [src_norm]
    for rv in raw_values:
        texts.append(normalize_text(rv))
        for sf in catalog.raw_to_sources.get(rv, set()):
            texts.append(normalize_text(sf))

    # Generic guardrail: preop variables should not use sources explicitly marked postoperative.
    if phase == "preop":
        for t in texts:
            if "postop" in t or "postoperative" in t or "postdischarge" in t:
                return True

    key = (normalize_text(var_name), db_key)
    blocked = FORBIDDEN_MATCH_PATTERNS.get(key, [])
    if not blocked:
        return False
    for t in texts:
        for pat in blocked:
            if pat in t:
                return True
    return False


def preop_time_check(matched_raw: Sequence[str], catalog: DBCatalog) -> Tuple[str, str]:
    if not matched_raw:
        return "NOT_CHECKED", ""

    texts: List[str] = []
    for raw in matched_raw:
        texts.append(normalize_text(raw))
        for sf in catalog.raw_to_sources.get(raw, set()):
            texts.append(normalize_text(sf))

    for t in texts:
        if any(tok in t for tok in PREOP_NEGATIVE_TOKENS):
            return "FAIL", f"negative token in matched evidence: {t}"
        # Common postoperative window patterns (e.g., 0_24, 0_48, 0_7d).
        if re.search(r"\b0\s*(24|48|72)\b", t) or re.search(r"\b0\s*(7|14|30)\s*d\b", t):
            return "FAIL", f"post-op window pattern in matched evidence: {t}"

    for t in texts:
        if any(tok in t for tok in PREOP_POSITIVE_TOKENS):
            return "PASS", f"preop evidence: {t}"

    return "UNCERTAIN", "no explicit preop/baseline/admission token in matched field or source file"


def preop_comorbidity_check(matched_raw: Sequence[str], catalog: DBCatalog) -> Tuple[str, str]:
    if not matched_raw:
        return "NOT_CHECKED", ""

    texts: List[str] = []
    for raw in matched_raw:
        texts.append(normalize_text(raw))
        for sf in catalog.raw_to_sources.get(raw, set()):
            texts.append(normalize_text(sf))

    for t in texts:
        if any(tok in t for tok in PREOP_NEGATIVE_TOKENS):
            return "FAIL", f"negative peri/post-op token in matched evidence: {t}"
        if re.search(r"\bpost\s*op.*complication", t) or re.search(r"\bpostoperative.*complication", t):
            return "FAIL", f"post-op complication evidence instead of baseline history: {t}"
        if "outcome" in t:
            return "FAIL", f"outcome evidence instead of baseline history: {t}"
        if re.search(r"\bintra\s*(1|2|3)\b", t) or re.search(r"\bintra(1|2|3)\b", t):
            return "FAIL", f"intra-phase marker in matched evidence: {t}"

    for t in texts:
        if any(tok in t for tok in COMORBIDITY_POSITIVE_TOKENS):
            return "PASS", f"history/comorbidity evidence: {t}"

    # For chronic comorbidity fields without explicit timing token,
    # accept when no peri/post-op signal is found.
    return "PASS", "no peri/post-op token found; accepted as baseline comorbidity field"


def match_variable(var_name: str, db_key: str, catalog: DBCatalog, category: str = "") -> Tuple[bool, List[str]]:
    nv = normalize_text(var_name)
    if not nv:
        return False, []

    if nv in catalog.norm_to_raw:
        return True, sorted(catalog.norm_to_raw[nv])[:3]

    ctx_free = remove_context_words(nv)

    core_toks = derive_key_tokens(var_name)
    if not core_toks:
        core_toks = {nv}
    expected_phase = variable_phase(var_name)
    if category == "术中化验":
        # In this project, intraop labs are proxy-derived from preop analyte list.
        expected_phase = ""
    med_var_phrase_aliases: List[str] = []
    if category == "住院用药":
        med_var_phrase_aliases = MED_VAR_SYNONYM_PHRASES.get(nv, [])

    cand_ids: Set[int] = set()
    for t in core_toks:
        cand_ids.update(catalog.token_index.get(t, set()))
        for a in TOKEN_ALIAS.get(t, set()):
            cand_ids.update(catalog.token_index.get(a, set()))
    if not cand_ids:
        phrase_aliases = []
        for t in core_toks:
            phrase_aliases.extend(TOKEN_ALIAS_PHRASE.get(t, []))
        if phrase_aliases:
            for i, src_norm in enumerate(catalog.norm_sources):
                if any(phrase_in_source(ph, src_norm) for ph in phrase_aliases):
                    cand_ids.add(i)
    if med_var_phrase_aliases:
        for i, src_norm in enumerate(catalog.norm_sources):
            if any(phrase_in_source(ph, src_norm) for ph in med_var_phrase_aliases):
                cand_ids.add(i)

    best_coverage = 0.0
    best_matched = 0
    best_phase_score = -1
    best_ctx_score = -1
    best_specificity = -1
    best_raw: List[str] = []

    for i in cand_ids:
        src_norm = catalog.norm_sources[i]
        raw_values = sorted(catalog.norm_to_raw[src_norm])
        if is_forbidden_candidate(var_name, db_key, src_norm, raw_values, catalog):
            continue
        src_toks = catalog.source_tokens[i]
        sp = source_phase(src_norm)
        if expected_phase and sp and expected_phase != sp:
            continue

        phase_score = 1
        if expected_phase:
            phase_score = 2 if sp == expected_phase else 1

        matched = 0
        for t in core_toks:
            acceptable_tokens = {t} | TOKEN_ALIAS.get(t, set())
            phrase_hits = any(phrase_in_source(ph, src_norm) for ph in TOKEN_ALIAS_PHRASE.get(t, []))
            if acceptable_tokens & src_toks or phrase_hits:
                matched += 1

        med_phrase_hit = bool(med_var_phrase_aliases) and any(
            phrase_in_source(ph, src_norm) for ph in med_var_phrase_aliases
        )
        if med_phrase_hit:
            matched = max(matched, len(core_toks))

        if matched == 0:
            continue
        coverage = matched / float(len(core_toks))
        ctx_score = 1 if (ctx_free and ctx_free == src_norm) else 0
        specificity = len(src_toks)

        better = False
        if phase_score > best_phase_score:
            better = True
        elif phase_score == best_phase_score and coverage > best_coverage:
            better = True
        elif phase_score == best_phase_score and coverage == best_coverage and matched > best_matched:
            better = True
        elif phase_score == best_phase_score and coverage == best_coverage and matched == best_matched and ctx_score > best_ctx_score:
            better = True
        elif (
            phase_score == best_phase_score
            and coverage == best_coverage
            and matched == best_matched
            and ctx_score == best_ctx_score
            and specificity > best_specificity
        ):
            better = True

        if better:
            best_phase_score = phase_score
            best_coverage = coverage
            best_matched = matched
            best_ctx_score = ctx_score
            best_specificity = specificity
            best_raw = raw_values[:3]

    if len(core_toks) == 1:
        return (best_coverage >= 1.0, best_raw if best_coverage >= 1.0 else [])
    if len(core_toks) == 2:
        # For 2-token variables, one strong token can be enough for abbreviation-heavy schemas.
        ok = best_coverage >= 0.5
        return ok, (best_raw if ok else [])
    ok = best_coverage >= 0.67
    return ok, (best_raw if ok else [])


def build_alignment_rows(
    categories: Dict[str, List[str]],
    catalogs: Dict[str, Dict[str, DBCatalog]],
    selected_categories: Sequence[str],
) -> List[Dict[str, str]]:
    rows: List[Dict[str, str]] = []
    for cat in selected_categories:
        domain = CATEGORY_DOMAIN.get(cat, "headers")
        for var in categories.get(cat, []):
            row = {"category": cat, "variable": var, "domain": domain}
            for db in DB_KEYS:
                catalog = catalogs[db][domain]
                has, matched = override_match(var, db, catalog)
                if not has:
                    has, matched = match_variable(var, db, catalog, category=cat)
                row[f"{db}_time_check"] = "NOT_APPLICABLE"
                row[f"{db}_time_evidence"] = ""
                note = "matched" if has else f"not found in scanned {domain} sources"
                if cat in PREOP_TIME_CHECK_CATEGORIES and has:
                    if db == "TJCH" and cat == "住院用药" and normalize_text(var) in TJCH_PREOP_MED_OVERRIDES:
                        time_status, time_evidence = "PASS", "manual rule: user-confirmed TJCH preoperative medication field"
                    else:
                        time_status, time_evidence = preop_time_check(matched, catalog)
                    row[f"{db}_time_check"] = time_status
                    row[f"{db}_time_evidence"] = time_evidence
                    if time_status == "FAIL":
                        has = False
                        note = "matched but failed preop time check"
                    elif time_status == "UNCERTAIN":
                        has = False
                        note = "matched but preop time evidence is uncertain (downgraded to ×)"
                elif cat == "合并症" and has:
                    time_status, time_evidence = preop_comorbidity_check(matched, catalog)
                    row[f"{db}_time_check"] = time_status
                    row[f"{db}_time_evidence"] = time_evidence
                    if time_status == "FAIL":
                        has = False
                        note = "matched but failed preop-comorbidity time check"
                row[db] = "√" if has else "×"
                row[f"{db}_match"] = " | ".join(matched)
                row[f"{db}_source_files"] = " | ".join(catalog.sources_for_raw(matched, limit=10))
                row[f"{db}_note"] = note
            rows.append(row)
    return rows


def write_csv(path: Path, rows: List[Dict[str, str]]) -> None:
    if not rows:
        return
    fieldnames = list(rows[0].keys())
    with path.open("w", encoding="utf-8", newline="") as f:
        w = csv.DictWriter(f, fieldnames=fieldnames)
        w.writeheader()
        w.writerows(rows)


def write_summary_csv(path: Path, rows: List[Dict[str, str]], selected_categories: Sequence[str]) -> None:
    summary: List[Dict[str, str]] = []
    for cat in selected_categories:
        chunk = [r for r in rows if r["category"] == cat]
        if not chunk:
            row = {"category": cat, "n_variables": "0"}
            for db in DB_KEYS:
                row[f"{db}_n"] = "0"
            summary.append(row)
            continue
        row = {"category": cat, "n_variables": str(len(chunk))}
        for db in DB_KEYS:
            row[f"{db}_n"] = str(sum(r[db] == "√" for r in chunk))
        summary.append(row)
    with path.open("w", encoding="utf-8", newline="") as f:
        fieldnames = ["category", "n_variables"] + [f"{db}_n" for db in DB_KEYS]
        w = csv.DictWriter(f, fieldnames=fieldnames)
        w.writeheader()
        w.writerows(summary)


def write_long_trace_csv(path: Path, rows: List[Dict[str, str]]) -> None:
    long_rows: List[Dict[str, str]] = []
    for r in rows:
        for db in DB_KEYS:
            long_rows.append(
                {
                    "category": r["category"],
                    "domain": r.get("domain", ""),
                    "variable": r["variable"],
                    "database": db,
                    "status": r.get(db, ""),
                    "matched_fields": r.get(f"{db}_match", ""),
                    "source_files": r.get(f"{db}_source_files", ""),
                    "time_check": r.get(f"{db}_time_check", ""),
                    "time_evidence": r.get(f"{db}_time_evidence", ""),
                    "note": r.get(f"{db}_note", ""),
                }
            )
    with path.open("w", encoding="utf-8", newline="") as f:
        fieldnames = [
            "category",
            "domain",
            "variable",
            "database",
            "status",
            "matched_fields",
            "source_files",
            "time_check",
            "time_evidence",
            "note",
        ]
        w = csv.DictWriter(f, fieldnames=fieldnames)
        w.writeheader()
        w.writerows(long_rows)


def write_category_wide_csv(path: Path, rows: List[Dict[str, str]], category: str) -> None:
    chunk = [r for r in rows if r["category"] == category]
    if not chunk:
        with path.open("w", encoding="utf-8", newline="") as f:
            fieldnames = ["category", "variable"] + [
                f"{db}_{k}"
                for db in DB_KEYS
                for k in ["status", "match", "source_files", "time_check", "time_evidence", "note"]
            ]
            w = csv.DictWriter(f, fieldnames=fieldnames)
            w.writeheader()
        return
    out_rows: List[Dict[str, str]] = []
    for r in chunk:
        out = {"category": r["category"], "variable": r["variable"]}
        for db in DB_KEYS:
            out[f"{db}_status"] = r.get(db, "")
            out[f"{db}_match"] = r.get(f"{db}_match", "")
            out[f"{db}_source_files"] = r.get(f"{db}_source_files", "")
            out[f"{db}_time_check"] = r.get(f"{db}_time_check", "")
            out[f"{db}_time_evidence"] = r.get(f"{db}_time_evidence", "")
            out[f"{db}_note"] = r.get(f"{db}_note", "")
        out_rows.append(out)
    with path.open("w", encoding="utf-8", newline="") as f:
        fieldnames = list(out_rows[0].keys())
        w = csv.DictWriter(f, fieldnames=fieldnames)
        w.writeheader()
        w.writerows(out_rows)


def write_review_excel(path: Path, rows: List[Dict[str, str]]) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "variable_review"
    headers = ["category", "variable"]
    for db in DB_KEYS:
        headers.extend(
            [
                f"{db}_status",
                f"{db}_match",
                f"{db}_source_files",
                f"{db}_time_check",
                f"{db}_time_evidence",
                f"{db}_note",
            ]
        )
    ws.append(headers)
    for r in rows:
        vals = [r["category"], r["variable"]]
        for db in DB_KEYS:
            vals.extend(
                [
                    r.get(db, ""),
                    r.get(f"{db}_match", ""),
                    r.get(f"{db}_source_files", ""),
                    r.get(f"{db}_time_check", ""),
                    r.get(f"{db}_time_evidence", ""),
                    r.get(f"{db}_note", ""),
                ]
            )
        ws.append(vals)
    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions

    ws2 = wb.create_sheet("long_trace")
    headers2 = [
        "category",
        "domain",
        "variable",
        "database",
        "status",
        "matched_fields",
        "source_files",
        "time_check",
        "time_evidence",
        "note",
    ]
    ws2.append(headers2)
    for r in rows:
        for db in DB_KEYS:
            ws2.append(
                [
                    r["category"],
                    r.get("domain", ""),
                    r["variable"],
                    db,
                    r.get(db, ""),
                    r.get(f"{db}_match", ""),
                    r.get(f"{db}_source_files", ""),
                    r.get(f"{db}_time_check", ""),
                    r.get(f"{db}_time_evidence", ""),
                    r.get(f"{db}_note", ""),
                ]
            )
    ws2.freeze_panes = "A2"
    ws2.auto_filter.ref = ws2.dimensions
    wb.save(str(path))


def ratio_text(ok_n: int, total_n: int) -> str:
    if total_n <= 0:
        return "0/0 (0.0%)"
    pct = 100.0 * ok_n / float(total_n)
    return f"{ok_n}/{total_n} ({pct:.1f}%)"


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(12.95), Inches(6.8), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(120, 128, 140)


def add_page_header(prs: Presentation, slide, title: str, subtitle: str = "") -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.72))
    band.fill.solid()
    band.fill.fore_color.rgb = REPORT_COLORS["navy"]
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.12), Inches(6.4), Inches(0.32))
    ttf = title_box.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.runs[0]
    tr.font.bold = True
    tr.font.size = Pt(17)
    tr.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.47), Inches(6.5), Inches(0.22))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.alignment = PP_ALIGN.LEFT
        sr = sp.runs[0]
        sr.font.size = Pt(9)
        sr.font.color.rgb = RGBColor(230, 236, 248)


def add_cover_slide(prs: Presentation, timestamp: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 253)
    bg.line.fill.background()

    add_page_header(prs, slide, "INSPIRE / VitalDB / MOVER / TJCH 变量对齐评估", "正式汇报版（竖版）")

    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6.5), Inches(6.5))
    tf = box.text_frame
    tf.clear()
    lines = [
        "汇报目的",
        "基于 INSPIRE 标准变量定义，评估 VitalDB、MOVER、TJCH字段可对齐性。",
        "",
        "评估范围",
        "基本信息、术前化验、合并症、手术信息、术中生命体征、术中用药、住院用药、术中化验、术后化验、并发症。",
        "",
        "输出",
        "逐变量 √/× 清单、多页明细表、覆盖率汇总、复合主要结局构建建议。",
        "",
        "生成时间",
        timestamp,
    ]
    for i, line in enumerate(lines):
        if i == 0 or i == 3 or i == 6 or i == 9:
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            r = p.runs[0] if p.runs else p.add_run()
            r.font.bold = True
            r.font.size = Pt(16)
            r.font.color.rgb = REPORT_COLORS["blue"]
        else:
            p = tf.add_paragraph()
            p.text = line
            r = p.runs[0] if p.runs else p.add_run()
            r.font.size = Pt(13)
            r.font.color.rgb = REPORT_COLORS["text"]

    add_footer(slide, "Variable Alignment Report")


def add_text_slide(prs: Presentation, title: str, subtitle: str, lines: List[str], footer: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(prs, slide, title, subtitle)

    box = slide.shapes.add_textbox(Inches(0.45), Inches(1.05), Inches(6.6), Inches(11.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        r = p.runs[0] if p.runs else p.add_run()
        if line.startswith("步骤") or line.startswith("复合主要结局定义") or line.startswith("匹配规则"):
            r.font.bold = True
            r.font.size = Pt(13)
            r.font.color.rgb = REPORT_COLORS["blue"]
        else:
            r.font.size = Pt(11.5)
            r.font.color.rgb = REPORT_COLORS["text"]
    add_footer(slide, footer)


def format_table_cell(cell, text: str, font_size: int, bold: bool = False, align: int = PP_ALIGN.LEFT, color=None) -> None:
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color if color is not None else REPORT_COLORS["text"]


def add_summary_slide(prs: Presentation, rows: List[Dict[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(prs, slide, "覆盖率汇总", "按模块统计可匹配变量比例（√ / 总变量数）")

    table = slide.shapes.add_table(
        len(REQUESTED_CATEGORIES) + 1,
        2 + len(DB_KEYS),
        Inches(0.35),
        Inches(1.15),
        Inches(6.8),
        Inches(8.2),
    ).table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(0.6)
    for idx in range(len(DB_KEYS)):
        table.columns[idx + 2].width = Inches(1.1)

    headers = ["模块", "N"] + [DB_LABELS[db] for db in DB_KEYS]
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = REPORT_COLORS["navy"]
        format_table_cell(c, h, 11, bold=True, align=PP_ALIGN.CENTER, color=RGBColor(255, 255, 255))

    for i, cat in enumerate(REQUESTED_CATEGORIES, start=1):
        chunk = [r for r in rows if r["category"] == cat]
        n = len(chunk)
        vals = [cat, str(n)]
        for db in DB_KEYS:
            db_n = sum(r[db] == "√" for r in chunk)
            vals.append(ratio_text(db_n, n))
        for j, v in enumerate(vals):
            c = table.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = REPORT_COLORS["light_blue"] if i % 2 == 0 else RGBColor(255, 255, 255)
            format_table_cell(c, v, 10, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    add_footer(slide, "Summary by Category")


def brief_missing_variables(rows: List[Dict[str, str]], db: str, limit: int = 6) -> str:
    miss = [r["variable"] for r in rows if r.get(db) == "×"]
    if not miss:
        return "无"
    if len(miss) <= limit:
        return "；".join(miss)
    return "；".join(miss[:limit]) + f"；...共{len(miss)}项"


def add_module_summary_slide(prs: Presentation, category: str, rows: List[Dict[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(prs, slide, f"{category} 模块总结", "每模块单页核查：覆盖率 + 缺失变量")

    n = len(rows)
    table = slide.shapes.add_table(
        len(DB_KEYS) + 1,
        3,
        Inches(0.35),
        Inches(1.1),
        Inches(6.8),
        Inches(2.4),
    ).table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(1.6)
    table.columns[2].width = Inches(3.7)

    headers = ["数据库", "覆盖率", "缺失变量（节选）"]
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = REPORT_COLORS["navy"]
        format_table_cell(c, h, 11, bold=True, align=PP_ALIGN.CENTER, color=RGBColor(255, 255, 255))

    for i, db in enumerate(DB_KEYS, start=1):
        db_n = sum(r.get(db) == "√" for r in rows)
        vals = [DB_LABELS[db], ratio_text(db_n, n), brief_missing_variables(rows, db, limit=6)]
        for j, val in enumerate(vals):
            c = table.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = REPORT_COLORS["light_blue"] if i % 2 == 0 else RGBColor(255, 255, 255)
            if j == 0:
                format_table_cell(c, val, 10, bold=True, align=PP_ALIGN.CENTER)
            elif j == 1:
                format_table_cell(c, val, 10, align=PP_ALIGN.CENTER)
            else:
                format_table_cell(c, val, 9, align=PP_ALIGN.LEFT)

    box = slide.shapes.add_textbox(Inches(0.35), Inches(3.9), Inches(6.8), Inches(8.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"模块变量总数: {n}"
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(12)
    p.runs[0].font.color.rgb = REPORT_COLORS["blue"]

    for db in DB_KEYS:
        miss = [r["variable"] for r in rows if r.get(db) == "×"]
        p = tf.add_paragraph()
        p.text = f"{DB_LABELS[db]} 缺失 {len(miss)} 项: " + ("；".join(miss) if miss else "无")
        r = p.runs[0] if p.runs else p.add_run()
        r.font.size = Pt(9.5)
        r.font.color.rgb = REPORT_COLORS["text"]

    add_footer(slide, "Module Summary")


def add_composite_alignment_slide(prs: Presentation, rows: List[Dict[str, str]]) -> None:
    by_var = {}
    for r in rows:
        if r["category"] == "并发症":
            by_var[r["variable"]] = r

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(prs, slide, "复合主要结局：跨库可用性", "组件级（任一组件命中即判为该库可用于复合结局构建）")

    table = slide.shapes.add_table(
        len(COMPOSITE_COMPONENTS) + 1,
        1 + len(DB_KEYS),
        Inches(0.35),
        Inches(1.15),
        Inches(6.8),
        Inches(5.2),
    ).table
    table.columns[0].width = Inches(2.8)
    for i in range(len(DB_KEYS)):
        table.columns[i + 1].width = Inches(1.0)

    for j, h in enumerate(["复合组件"] + [DB_LABELS[db] for db in DB_KEYS]):
        c = table.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = REPORT_COLORS["navy"]
        format_table_cell(c, h, 10, bold=True, align=PP_ALIGN.CENTER, color=RGBColor(255, 255, 255))

    for i, (comp_name, comp_vars) in enumerate(COMPOSITE_COMPONENTS, start=1):
        has = {db: False for db in DB_KEYS}
        for v in comp_vars:
            rv = by_var.get(v)
            if not rv:
                continue
            for db in has.keys():
                has[db] = has[db] or (rv[db] == "√")

        row_vals = [comp_name] + [("√" if has[db] else "×") for db in DB_KEYS]
        for j, val in enumerate(row_vals):
            c = table.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = REPORT_COLORS["light_blue"] if i % 2 == 0 else RGBColor(255, 255, 255)
            if j == 0:
                format_table_cell(c, val, 9, align=PP_ALIGN.LEFT)
            else:
                col = REPORT_COLORS["ok"] if val == "√" else REPORT_COLORS["no"]
                format_table_cell(c, val, 12, bold=True, align=PP_ALIGN.CENTER, color=col)

    note = [
        "复合主要结局定义",
        "MCO = 1（发生）当且仅当任一复合组件在结局窗口内发生；否则 MCO = 0。",
        "建议窗口",
        "首选住院期 + 术后30天；若某数据库仅有住院窗口，则按住院窗口构建并在模型层做窗口敏感性分析。",
    ]
    add_text_slide(
        prs,
        "复合主要结局：定义建议",
        "用于后续跨库统一建模",
        note,
        "Composite Outcome Definition",
    )

    add_footer(slide, "Composite Outcome Feasibility")


def add_detail_table_slide(prs: Presentation, title: str, table_rows: List[Dict[str, str]], page_label: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_header(prs, slide, f"{title}  ({page_label})", "变量可用性检查：√=可匹配，×=未匹配")

    n_rows = len(table_rows) + 1
    table = slide.shapes.add_table(
        n_rows,
        1 + len(DB_KEYS),
        Inches(0.35),
        Inches(1.05),
        Inches(6.8),
        Inches(11.6),
    ).table
    table.columns[0].width = Inches(3.9)
    for i in range(len(DB_KEYS)):
        table.columns[i + 1].width = Inches(0.72)

    for j, h in enumerate(["Variable"] + [DB_LABELS[db] for db in DB_KEYS]):
        c = table.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = REPORT_COLORS["navy"]
        format_table_cell(c, h, 10, bold=True, align=PP_ALIGN.CENTER, color=RGBColor(255, 255, 255))

    for i, r in enumerate(table_rows, start=1):
        vals = [r["variable"]] + [r[db] for db in DB_KEYS]
        for j, v in enumerate(vals):
            c = table.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = REPORT_COLORS["light_blue"] if i % 2 == 0 else RGBColor(255, 255, 255)
            if j == 0:
                format_table_cell(c, v, 8, align=PP_ALIGN.LEFT)
            else:
                col = REPORT_COLORS["ok"] if v == "√" else REPORT_COLORS["no"]
                format_table_cell(c, v, 11, bold=True, align=PP_ALIGN.CENTER, color=col)

    add_footer(slide, "Detailed Variable Checklist")


def write_ppt(
    path: Path,
    rows: List[Dict[str, str]],
    timestamp: str,
    selected_categories: Sequence[str],
    mode: str = "detail_only",
    detail_max_rows: int = 45,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(13.33)

    if mode == "full":
        add_cover_slide(prs, timestamp)
        add_text_slide(
            prs,
            "方法步骤",
            "数据标准、扫描、匹配、输出",
            [
                "步骤1：标准定义抽取",
                "从 Variable_definition_4_22_2026.docx 抽取变量并按10个模块组织。",
                "步骤2：四库字段采集",
                "INSPIRE 扫描全量 CSV 表头；VitalDB 扫描 processed 表头并抽取 lab name；MOVER 扫描 SIS/EPIC 表头并抽取 lab/med/dx/complication 值域；TJCH扫描 XLSX 表头。",
                "步骤3：域内匹配",
                "按模块限定匹配域（headers / labs / meds / dx / complications），先做规范化精确匹配，再做关键词与同义词规则匹配。",
                "步骤4：可用性判定",
                "每变量在每库给出 √ 或 ×，并保留匹配字段与原始来源文件用于人工复核。",
                "步骤5：报告输出",
                "输出明细 CSV、逐变量来源追踪 CSV、模块化单页总结 PPT。",
            ],
            "Workflow Overview",
        )
        add_text_slide(
            prs,
            "匹配规则",
            "用于跨库字段命名不一致的自动对齐",
            [
                "匹配规则",
                "1) 文本规范化：大小写、下划线、特殊字符、下标字符统一。",
                "2) 语境词剥离：去除 pre/intra/postoperative 等上下文词后再匹配。",
                "3) 同义词扩展：如 hb~hgb~hemoglobin, na~sodium, cr~creatinine, spo2~sp02。",
                "4) 域内限制：化验只对 lab 字典匹配，用药只对 med 字典匹配，避免跨域误匹配。",
                "5) 阈值：单关键词需完整命中；多关键词需达到覆盖阈值。",
                "说明：本报告用于“变量是否存在”的可用性筛查，不等同于临床语义完全一致校验。",
            ],
            "Matching Strategy",
        )
        add_summary_slide(prs, rows)
        add_composite_alignment_slide(prs, rows)
        for cat in selected_categories:
            chunk = [r for r in rows if r["category"] == cat]
            add_module_summary_slide(prs, cat, chunk)
    else:
        # Review mode: only keep old-style detail checklist slides by selected module(s).
        per_slide = max(10, int(detail_max_rows))
        for cat in selected_categories:
            chunk = [r for r in rows if r["category"] == cat]
            if not chunk:
                add_detail_table_slide(prs, cat, [], "1/1")
                continue
            n_pages = (len(chunk) + per_slide - 1) // per_slide
            for p in range(n_pages):
                st = p * per_slide
                ed = min((p + 1) * per_slide, len(chunk))
                add_detail_table_slide(prs, cat, chunk[st:ed], f"{p + 1}/{n_pages}")

    prs.save(str(path))


def write_notes(path: Path, docx_path: Path) -> None:
    text = f"""Alignment notes
================
1) Standard source: {docx_path}
2) Category mapping:
   - Basic info = Table 1 + Table 3.
   - In-hospital medication = Table 5.
   - Intraop vitals and intraop meds are split from Table 7 by keyword rules.
   - Intraop labs are derived from Table 4 analyte names (this DOCX has no standalone intraop-lab table).
3) Matching method:
   - normalized exact match
   - domain-constrained keyword/synonym match
   - value dictionaries for long tables (lab/medication/diagnosis/complication)
4) Evidence:
   - keep matched raw field names
   - keep raw source file/sheet paths for each matched field
5) Report scope: availability screen (√/×), not a full semantic harmonization audit.
"""
    path.write_text(text, encoding="utf-8")


def write_method_description(path: Path, rows: List[Dict[str, str]], timestamp: str) -> None:
    comp_rows = [r for r in rows if r["category"] == "并发症"]
    by_var = {r["variable"]: r for r in comp_rows}

    lines = []
    lines.append("# 四库变量对齐与复合主要结局构建说明")
    lines.append("")
    lines.append("## 1) 目标")
    lines.append("以 INSPIRE 变量定义为标准，检查 VitalDB、MOVER、TJCH是否存在可对应字段，形成跨库可用性矩阵。")
    lines.append("")
    lines.append("## 2) 数据范围")
    lines.append("- 标准文档：Variable_definition_4_22_2026.docx")
    lines.append("- 数据根目录：/N/project/analgesia_perioperation/data")
    lines.append("- 四个数据库：INSPIRE_1.3、VitalDB_1.0.0、MOVER、TJCH（cohort_data_10_10_2025_final.xlsx）")
    lines.append("")
    lines.append("## 3) 对齐步骤")
    lines.append("1. 从标准文档抽取变量，按10个模块组织。")
    lines.append("2. 采集各库字段与值域字典（lab/med/dx/complication）。")
    lines.append("3. 在模块对应域内进行匹配：精确匹配 + 关键词/同义词匹配。")
    lines.append("4. 输出逐变量 √/× 结果、匹配字段、原始来源文件路径。")
    lines.append("5. 生成模块化单页总结 PPT 与核查 CSV。")
    lines.append("")
    lines.append("## 4) 复合主要结局（MCO）建议定义")
    lines.append("定义：若患者在结局窗口内发生任一复合组件，则 MCO=1；否则 MCO=0。")
    lines.append("建议窗口：住院期 + 术后30天；如库内仅有住院窗口，则按住院窗口构建并在分析时做敏感性分析。")
    lines.append("")
    lines.append("复合组件建议：")
    for comp_name, comp_vars in COMPOSITE_COMPONENTS:
        lines.append(f"- {comp_name}: " + "；".join(comp_vars))
    lines.append("")
    lines.append("## 5) 复合组件跨库可用性（本次自动匹配）")
    lines.append("| 组件 | INSPIRE | VitalDB | MOVER | TJCH |")
    lines.append("|---|---|---|---|---|")
    for comp_name, comp_vars in COMPOSITE_COMPONENTS:
        has = {db: False for db in DB_KEYS}
        for v in comp_vars:
            rv = by_var.get(v)
            if not rv:
                continue
            for db in has.keys():
                has[db] = has[db] or (rv[db] == "√")
        lines.append(
            f"| {comp_name} | {'√' if has['INSPIRE'] else '×'} | {'√' if has['VitalDB'] else '×'} | {'√' if has['MOVER'] else '×'} | {'√' if has['TJCH'] else '×'} |"
        )
    lines.append("")
    lines.append("## 6) 注意事项")
    lines.append("- 本结果是“变量存在性/可检索性”筛查，不代表定义完全同构。")
    lines.append("- 进入建模前，需要对时间窗口、单位、阈值、重复记录汇总规则做语义对齐。")
    lines.append("")
    lines.append("生成时间: " + timestamp)

    path.write_text("\n".join(lines), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser()
    p.add_argument(
        "--docx",
        type=Path,
        default=Path("/N/project/analgesia_perioperation/projects/Inspire_data_process_ZZ/Variable_definition_4_22_2026.docx"),
    )
    p.add_argument(
        "--data-root",
        type=Path,
        default=Path("/N/project/analgesia_perioperation/data"),
    )
    p.add_argument(
        "--out-dir",
        type=Path,
        default=Path("/N/project/analgesia_perioperation/projects/Inspire_data_process_ZZ/06_Results/variable_alignment_20260422"),
    )
    p.add_argument(
        "--tjth-xlsx",
        type=Path,
        default=Path("/N/project/analgesia_perioperation/projects/Inspire_data_process_ZZ/cohort_data_10_10_2025_final.xlsx"),
    )
    p.add_argument(
        "--modules",
        type=str,
        default="基本信息",
        help="Comma-separated categories to review. Example: 基本信息,术前化验",
    )
    p.add_argument(
        "--ppt-mode",
        type=str,
        choices=["detail_only", "full"],
        default="detail_only",
        help="detail_only: old-style variable checklist tables only; full: full deck",
    )
    p.add_argument(
        "--detail-max-rows",
        type=int,
        default=45,
        help="Max table rows per slide in detail_only mode",
    )
    return p.parse_args()


def main() -> None:
    args = parse_args()
    args.out_dir.mkdir(parents=True, exist_ok=True)
    selected_categories = parse_modules_arg(args.modules)

    categories = extract_standard_variables(args.docx)
    catalogs = build_catalogs(args.data_root, args.tjth_xlsx)
    rows = build_alignment_rows(categories, catalogs, selected_categories)

    detail_csv = args.out_dir / "variable_alignment_detail.csv"
    summary_csv = args.out_dir / "variable_alignment_summary.csv"
    trace_long_csv = args.out_dir / "variable_alignment_trace_long.csv"
    if len(selected_categories) == 1:
        review_slug = CATEGORY_SLUG.get(selected_categories[0], "module")
    else:
        review_slug = "selected_modules"
    review_xlsx = args.out_dir / f"{review_slug}_alignment_review.xlsx"
    pptx_path = args.out_dir / "variable_alignment_checklist.pptx"
    notes_path = args.out_dir / "README_alignment_notes.txt"
    method_path = args.out_dir / "method_description_alignment_and_composite.md"

    write_csv(detail_csv, rows)
    write_summary_csv(summary_csv, rows, selected_categories)
    write_long_trace_csv(trace_long_csv, rows)
    for cat in selected_categories:
        cat_slug = CATEGORY_SLUG.get(cat, "category")
        cat_csv = args.out_dir / f"{cat_slug}_variable_source_trace.csv"
        write_category_wide_csv(cat_csv, rows, cat)
    write_review_excel(review_xlsx, rows)
    write_ppt(
        pptx_path,
        rows,
        datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        selected_categories=selected_categories,
        mode=args.ppt_mode,
        detail_max_rows=args.detail_max_rows,
    )
    write_notes(notes_path, args.docx)
    write_method_description(method_path, rows, datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

    print(f"Wrote: {detail_csv}")
    print(f"Wrote: {summary_csv}")
    print(f"Wrote: {trace_long_csv}")
    for cat in selected_categories:
        print(f"Wrote: {args.out_dir / (CATEGORY_SLUG.get(cat, 'category') + '_variable_source_trace.csv')}")
    print(f"Wrote: {review_xlsx}")
    print(f"Wrote: {pptx_path}")
    print(f"Wrote: {notes_path}")
    print(f"Wrote: {method_path}")


if __name__ == "__main__":
    main()
